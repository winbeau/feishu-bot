import csv
import os
from pathlib import Path
from statistics import mean

from app.core.models import Attachment


class FileParserService:
    def __init__(
        self,
        full_text_max_bytes: int | None = None,
        parsed_text_max_chars: int | None = None,
    ) -> None:
        self._full_text_max_bytes = int(
            full_text_max_bytes
            if full_text_max_bytes is not None
            else os.getenv("FILE_FULL_TEXT_MAX_BYTES", "65536")
        )
        self._parsed_text_max_chars = int(
            parsed_text_max_chars
            if parsed_text_max_chars is not None
            else os.getenv("DIFY_PARSED_TEXT_MAX_CHARS", "12000")
        )

    def parse_attachment(self, attachment: Attachment) -> Attachment:
        try:
            if not attachment.local_path:
                attachment.file_tags.append("parse_skipped:no_local_path")
                return attachment

            path = Path(attachment.local_path)
            suffix = path.suffix.lower().lstrip(".")
            mime_type = (attachment.mime_type or "").lower()
            if suffix in {"md", "txt"} or mime_type.startswith("text/"):
                text = self._parse_text(path, attachment)
            elif suffix == "csv" or mime_type in {"text/csv", "application/csv"}:
                text = self._parse_csv(path, attachment)
            elif suffix == "pptx" or mime_type.endswith(
                "presentationml.presentation"
            ):
                text = self._parse_pptx(path)
            elif suffix == "pdf" or mime_type == "application/pdf":
                text = self._parse_pdf(path, attachment)
            else:
                text = self._parse_text(path, attachment)
                attachment.file_tags.append("parsed_as_text")

            attachment.parsed_text = self._truncate_chars(text, attachment)
        except Exception as exc:
            attachment.parsed_text = None
            attachment.file_tags.extend(["parse_failed", exc.__class__.__name__])
        return attachment

    def _parse_text(self, path: Path, attachment: Attachment) -> str:
        data = path.read_bytes()
        if len(data) > self._full_text_max_bytes:
            data = data[: self._full_text_max_bytes]
            attachment.file_tags.append("truncated_bytes")
        return data.decode("utf-8", errors="replace")

    def _parse_csv(self, path: Path, attachment: Attachment) -> str:
        if path.stat().st_size <= self._full_text_max_bytes:
            rows = self._read_csv_rows(path)
            return self._markdown_table(rows)

        rows = self._read_csv_rows(path)
        if not rows:
            return ""

        header = rows[0]
        body = rows[1:]
        numeric_stats = self._numeric_stats(header, body)
        parts = [
            f"Columns: {', '.join(header)}",
            f"Rows: {len(body)}",
            "First 5 rows:",
            self._markdown_table([header, *body[:5]]),
        ]
        if numeric_stats:
            parts.extend(["Numeric summary:", *numeric_stats])
        attachment.file_tags.append("csv_summary_only")
        return "\n\n".join(parts)

    def _read_csv_rows(self, path: Path) -> list[list[str]]:
        with path.open(newline="", encoding="utf-8", errors="replace") as handle:
            return [row for row in csv.reader(handle)]

    def _markdown_table(self, rows: list[list[str]]) -> str:
        if not rows:
            return ""
        width = max(len(row) for row in rows)
        normalized = [row + [""] * (width - len(row)) for row in rows]
        header = normalized[0]
        separator = ["---"] * width
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in normalized[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)

    def _numeric_stats(self, header: list[str], rows: list[list[str]]) -> list[str]:
        stats: list[str] = []
        for index, column_name in enumerate(header):
            values: list[float] = []
            for row in rows:
                if index >= len(row) or row[index] == "":
                    continue
                try:
                    values.append(float(row[index]))
                except ValueError:
                    continue
            if values:
                stats.append(
                    f"- {column_name}: count={len(values)}, min={min(values)}, "
                    f"max={max(values)}, mean={mean(values):.2f}"
                )
        return stats

    def _parse_pptx(self, path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(str(path))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            slides.append(f"Slide {index}\n" + "\n".join(texts))
        return "\n\n".join(slides)

    def _parse_pdf(self, path: Path, attachment: Attachment) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        if len(pages) > 1:
            attachment.file_tags.append(f"pdf_pages:{len(pages)}")
        return "\n\n".join(pages)

    def _truncate_chars(self, text: str, attachment: Attachment) -> str:
        if len(text) <= self._parsed_text_max_chars:
            return text
        attachment.file_tags.append("truncated_chars")
        return text[: self._parsed_text_max_chars]
