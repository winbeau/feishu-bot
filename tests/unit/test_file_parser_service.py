from pathlib import Path

from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from app.core.models import Attachment
from app.services.file_parser import FileParserService


def write_text(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def test_file_parser_reads_text_and_truncates_long_content(tmp_path) -> None:
    path = tmp_path / "note.txt"
    write_text(path, "abcdef")
    service = FileParserService(full_text_max_bytes=3, parsed_text_max_chars=10)
    attachment = Attachment(local_path=str(path), mime_type="text/plain")

    parsed = service.parse_attachment(attachment)

    assert parsed.parsed_text == "abc"
    assert "truncated_bytes" in parsed.file_tags


def test_file_parser_small_csv_outputs_markdown_table(tmp_path) -> None:
    path = tmp_path / "data.csv"
    write_text(path, "name,value\nalpha,1\nbeta,2\n")
    service = FileParserService(full_text_max_bytes=1000)

    parsed = service.parse_attachment(Attachment(local_path=str(path)))

    assert parsed.parsed_text == (
        "| name | value |\n"
        "| --- | --- |\n"
        "| alpha | 1 |\n"
        "| beta | 2 |"
    )


def test_file_parser_large_csv_outputs_summary_not_full_text(tmp_path) -> None:
    path = tmp_path / "data.csv"
    write_text(path, "name,value\nalpha,1\nbeta,2\ngamma,3\ndelta,4\nepsilon,5\nzeta,6\n")
    service = FileParserService(full_text_max_bytes=10)

    parsed = service.parse_attachment(Attachment(local_path=str(path)))

    assert "Columns: name, value" in (parsed.parsed_text or "")
    assert "Rows: 6" in (parsed.parsed_text or "")
    assert "mean=3.50" in (parsed.parsed_text or "")
    assert "csv_summary_only" in parsed.file_tags


def test_file_parser_extracts_pptx_slide_text(tmp_path) -> None:
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Revenue up"
    presentation.save(path)
    service = FileParserService()

    parsed = service.parse_attachment(Attachment(local_path=str(path)))

    assert "Slide 1" in (parsed.parsed_text or "")
    assert "Quarterly Review" in (parsed.parsed_text or "")
    assert "Revenue up" in (parsed.parsed_text or "")


def test_file_parser_extracts_pdf_text_and_page_tag(tmp_path) -> None:
    path = tmp_path / "doc.pdf"
    writer = PdfWriter()
    page = writer.add_blank_page(width=300, height=144)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )
    stream = DecodedStreamObject()
    stream.set_data(b"BT /F1 24 Tf 72 72 Td (Hello PDF) Tj ET")
    page[NameObject("/Contents")] = writer._add_object(stream)
    with path.open("wb") as handle:
        writer.write(handle)
    service = FileParserService()

    parsed = service.parse_attachment(Attachment(local_path=str(path)))

    assert "Hello PDF" in (parsed.parsed_text or "")


def test_file_parser_marks_parse_failed_without_raising(tmp_path) -> None:
    path = tmp_path / "missing.txt"
    service = FileParserService()

    parsed = service.parse_attachment(Attachment(local_path=str(path)))

    assert parsed.parsed_text is None
    assert parsed.file_tags[0] == "parse_failed"
